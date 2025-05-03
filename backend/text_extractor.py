import pdfplumber
from docx import Document
from pptx import Presentation
import os

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(docx_path):
    """Extract text from a DOCX file."""
    doc = Document(docx_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    """Extract text from a PPTX file."""
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path):
    """Extract text from any supported file type."""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}") 